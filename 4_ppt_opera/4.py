import pptx
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


p = pptx.Presentation()

layout = p.slide_layouts[1]
slide = p.slides.add_slide(layout)

title = slide.placeholders[0]
title.text = '题目'
content = slide.placeholders[1]

# 自定义段落
paragraph_1 = content.text_frame.add_paragraph()
paragraph_1.text = '第一段'
paragraph_1.bold = True
paragraph_1.font.italic = True
paragraph_1.font.size = Pt(16)
paragraph_1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

paragraph_2 = content.text_frame.add_paragraph()
paragraph_2.text = '第二段'
paragraph_2.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT


# 插入表格
layout_two = p.slide_layouts[1]
slide_two = p.slides.add_slide(layout)
rows = 6
cols = 2

left = top = Inches(1)
width = Inches(6.0)
height = Inches(1.0)

table = slide_two.shapes.add_table(rows, cols, left, top, width, height).table

for r in range(rows):
    for c in range(cols):
        table.cell(r, c).text = f'{r}-{c}'


# 插入图片
layout_three = p.slide_layouts[1]
slide_three = p.slides.add_slide(layout)

left = top = Inches(1)
width = Inches(6.0)
height = Inches(6.0)
image = slide_three.shapes.add_picture(image_file='1.jpeg', left=left, top=top, width=width, height=height)

p.save('test.ppt')


