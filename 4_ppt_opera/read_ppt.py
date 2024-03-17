import pptx

p = pptx.Presentation('test.ppt')
for slide in p.slides:
    # 每一页有很多shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text_frame.text)
        if shape.has_table:
            for cell in shape.table.iter_cells():
                print(cell.text)
